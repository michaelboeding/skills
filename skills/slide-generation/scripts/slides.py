#!/usr/bin/env python3
"""
Slide Generation Script

Creates presentation slides from structured JSON content.
Outputs to PPTX (PowerPoint) or Markdown format.

Usage:
  python slides.py --content slides.json --output presentation.pptx
  python slides.py --content slides.json --theme modern-dark --output deck.pptx
  python slides.py --content slides.json --format markdown --output slides.md

Requirements:
  pip install python-pptx Pillow
"""

import argparse
import json
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RgbColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

# Theme definitions
THEMES = {
    "modern-dark": {
        "background": "1a1a2e",
        "text": "ffffff",
        "accent": "e94560",
        "secondary": "0f3460",
        "heading_size": 36,
        "body_size": 18,
    },
    "modern-light": {
        "background": "ffffff",
        "text": "1f2937",
        "accent": "3b82f6",
        "secondary": "e5e7eb",
        "heading_size": 36,
        "body_size": 18,
    },
    "corporate": {
        "background": "ffffff",
        "text": "1e3a5f",
        "accent": "2563eb",
        "secondary": "dbeafe",
        "heading_size": 32,
        "body_size": 16,
    },
    "startup": {
        "background": "0f172a",
        "text": "f8fafc",
        "accent": "f59e0b",
        "secondary": "1e293b",
        "heading_size": 40,
        "body_size": 20,
    },
    "minimal": {
        "background": "fafafa",
        "text": "18181b",
        "accent": "18181b",
        "secondary": "e4e4e7",
        "heading_size": 44,
        "body_size": 18,
    },
}


def hex_to_rgb(hex_color: str) -> tuple:
    """Convert hex color to RGB tuple."""
    hex_color = hex_color.lstrip('#')
    return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))


def set_slide_background(slide, hex_color: str):
    """Set slide background color."""
    background = slide.background
    fill = background.fill
    fill.solid()
    r, g, b = hex_to_rgb(hex_color)
    fill.fore_color.rgb = RgbColor(r, g, b)


def add_title_slide(prs, slide_data: dict, theme: dict):
    """Create a title slide."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, theme["background"])
    
    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(2.5), Inches(9), Inches(1.5)
    )
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = slide_data.get("title", "")
    title_para.font.size = Pt(theme["heading_size"] + 8)
    title_para.font.bold = True
    r, g, b = hex_to_rgb(theme["text"])
    title_para.font.color.rgb = RgbColor(r, g, b)
    title_para.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    if slide_data.get("subtitle"):
        sub_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(4), Inches(9), Inches(1)
        )
        sub_frame = sub_box.text_frame
        sub_para = sub_frame.paragraphs[0]
        sub_para.text = slide_data["subtitle"]
        sub_para.font.size = Pt(theme["body_size"] + 4)
        r, g, b = hex_to_rgb(theme["accent"])
        sub_para.font.color.rgb = RgbColor(r, g, b)
        sub_para.alignment = PP_ALIGN.CENTER
    
    return slide


def add_content_slide(prs, slide_data: dict, theme: dict):
    """Create a content slide with bullets."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, theme["background"])
    
    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.5), Inches(9), Inches(1)
    )
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = slide_data.get("title", "")
    title_para.font.size = Pt(theme["heading_size"])
    title_para.font.bold = True
    r, g, b = hex_to_rgb(theme["text"])
    title_para.font.color.rgb = RgbColor(r, g, b)
    
    # Bullets
    bullets = slide_data.get("bullets", [])
    if bullets:
        content_box = slide.shapes.add_textbox(
            Inches(0.75), Inches(1.75), Inches(8.5), Inches(5)
        )
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        
        for i, bullet in enumerate(bullets):
            if i == 0:
                para = content_frame.paragraphs[0]
            else:
                para = content_frame.add_paragraph()
            para.text = f"• {bullet}"
            para.font.size = Pt(theme["body_size"])
            r, g, b = hex_to_rgb(theme["text"])
            para.font.color.rgb = RgbColor(r, g, b)
            para.space_after = Pt(12)
    
    return slide


def add_section_slide(prs, slide_data: dict, theme: dict):
    """Create a section divider slide."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Use accent color for background
    set_slide_background(slide, theme["accent"])
    
    # Section title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(3), Inches(9), Inches(1.5)
    )
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = slide_data.get("title", "")
    title_para.font.size = Pt(theme["heading_size"] + 4)
    title_para.font.bold = True
    # Use background color for text (inverted)
    r, g, b = hex_to_rgb(theme["background"])
    title_para.font.color.rgb = RgbColor(r, g, b)
    title_para.alignment = PP_ALIGN.CENTER
    
    return slide


def add_quote_slide(prs, slide_data: dict, theme: dict):
    """Create a quote slide."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, theme["background"])
    
    # Quote
    quote_box = slide.shapes.add_textbox(
        Inches(1), Inches(2), Inches(8), Inches(3)
    )
    quote_frame = quote_box.text_frame
    quote_para = quote_frame.paragraphs[0]
    quote_para.text = f'"{slide_data.get("quote", "")}"'
    quote_para.font.size = Pt(theme["body_size"] + 6)
    quote_para.font.italic = True
    r, g, b = hex_to_rgb(theme["text"])
    quote_para.font.color.rgb = RgbColor(r, g, b)
    quote_para.alignment = PP_ALIGN.CENTER
    
    # Attribution
    if slide_data.get("attribution"):
        attr_box = slide.shapes.add_textbox(
            Inches(1), Inches(5), Inches(8), Inches(0.5)
        )
        attr_frame = attr_box.text_frame
        attr_para = attr_frame.paragraphs[0]
        attr_para.text = slide_data["attribution"]
        attr_para.font.size = Pt(theme["body_size"])
        r, g, b = hex_to_rgb(theme["accent"])
        attr_para.font.color.rgb = RgbColor(r, g, b)
        attr_para.alignment = PP_ALIGN.CENTER
    
    return slide


def add_stats_slide(prs, slide_data: dict, theme: dict):
    """Create a stats slide."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, theme["background"])
    
    # Title
    if slide_data.get("title"):
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5), Inches(9), Inches(1)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = slide_data["title"]
        title_para.font.size = Pt(theme["heading_size"])
        title_para.font.bold = True
        r, g, b = hex_to_rgb(theme["text"])
        title_para.font.color.rgb = RgbColor(r, g, b)
    
    # Stats
    stats = slide_data.get("stats", [])
    if stats:
        num_stats = len(stats)
        stat_width = 8 / num_stats
        start_x = 1
        
        for i, stat in enumerate(stats):
            x = start_x + (i * stat_width)
            
            # Value
            value_box = slide.shapes.add_textbox(
                Inches(x), Inches(2.5), Inches(stat_width - 0.2), Inches(1.5)
            )
            value_frame = value_box.text_frame
            value_para = value_frame.paragraphs[0]
            value_para.text = str(stat.get("value", ""))
            value_para.font.size = Pt(48)
            value_para.font.bold = True
            r, g, b = hex_to_rgb(theme["accent"])
            value_para.font.color.rgb = RgbColor(r, g, b)
            value_para.alignment = PP_ALIGN.CENTER
            
            # Label
            label_box = slide.shapes.add_textbox(
                Inches(x), Inches(4), Inches(stat_width - 0.2), Inches(0.5)
            )
            label_frame = label_box.text_frame
            label_para = label_frame.paragraphs[0]
            label_para.text = stat.get("label", "")
            label_para.font.size = Pt(theme["body_size"])
            r, g, b = hex_to_rgb(theme["text"])
            label_para.font.color.rgb = RgbColor(r, g, b)
            label_para.alignment = PP_ALIGN.CENTER
    
    return slide


def add_closing_slide(prs, slide_data: dict, theme: dict):
    """Create a closing slide."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, theme["background"])
    
    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(2.5), Inches(9), Inches(1.5)
    )
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = slide_data.get("title", "Thank You")
    title_para.font.size = Pt(theme["heading_size"] + 8)
    title_para.font.bold = True
    r, g, b = hex_to_rgb(theme["text"])
    title_para.font.color.rgb = RgbColor(r, g, b)
    title_para.alignment = PP_ALIGN.CENTER
    
    # Subtitle/Contact
    if slide_data.get("subtitle") or slide_data.get("contact"):
        text = slide_data.get("subtitle", "") or slide_data.get("contact", "")
        sub_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(4), Inches(9), Inches(1)
        )
        sub_frame = sub_box.text_frame
        sub_para = sub_frame.paragraphs[0]
        sub_para.text = text
        sub_para.font.size = Pt(theme["body_size"])
        r, g, b = hex_to_rgb(theme["accent"])
        sub_para.font.color.rgb = RgbColor(r, g, b)
        sub_para.alignment = PP_ALIGN.CENTER
    
    return slide


def create_pptx(content: dict, theme_name: str, output_path: str):
    """Create a PowerPoint presentation."""
    if not HAS_PPTX:
        print("Error: python-pptx not installed. Run: pip install python-pptx", file=sys.stderr)
        sys.exit(1)
    
    # Get theme
    theme = THEMES.get(theme_name, THEMES["modern-dark"])
    if isinstance(content.get("metadata", {}).get("theme"), dict):
        theme.update(content["metadata"]["theme"])
    
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Process slides
    slides_data = content.get("slides", [])
    
    for slide_data in slides_data:
        slide_type = slide_data.get("type", "content")
        
        if slide_type == "title":
            add_title_slide(prs, slide_data, theme)
        elif slide_type == "content":
            add_content_slide(prs, slide_data, theme)
        elif slide_type == "section":
            add_section_slide(prs, slide_data, theme)
        elif slide_type == "quote":
            add_quote_slide(prs, slide_data, theme)
        elif slide_type == "stats":
            add_stats_slide(prs, slide_data, theme)
        elif slide_type == "closing":
            add_closing_slide(prs, slide_data, theme)
        else:
            # Default to content slide
            add_content_slide(prs, slide_data, theme)
    
    # Save
    prs.save(output_path)
    return len(slides_data)


def create_markdown(content: dict, output_path: str):
    """Create Marp-compatible Markdown slides."""
    lines = [
        "---",
        "marp: true",
        f"title: {content.get('metadata', {}).get('title', 'Presentation')}",
        "theme: default",
        "paginate: true",
        "---",
        ""
    ]
    
    slides_data = content.get("slides", [])
    
    for i, slide_data in enumerate(slides_data):
        slide_type = slide_data.get("type", "content")
        
        if i > 0:
            lines.append("---")
            lines.append("")
        
        if slide_type == "title":
            lines.append(f"# {slide_data.get('title', '')}")
            if slide_data.get("subtitle"):
                lines.append(f"## {slide_data['subtitle']}")
        
        elif slide_type == "content":
            lines.append(f"# {slide_data.get('title', '')}")
            lines.append("")
            for bullet in slide_data.get("bullets", []):
                lines.append(f"- {bullet}")
        
        elif slide_type == "section":
            lines.append(f"# {slide_data.get('title', '')}")
        
        elif slide_type == "quote":
            lines.append(f"> {slide_data.get('quote', '')}")
            if slide_data.get("attribution"):
                lines.append(f"> {slide_data['attribution']}")
        
        elif slide_type == "stats":
            lines.append(f"# {slide_data.get('title', '')}")
            lines.append("")
            stats = slide_data.get("stats", [])
            for stat in stats:
                lines.append(f"**{stat.get('value', '')}** {stat.get('label', '')}")
        
        elif slide_type == "closing":
            lines.append(f"# {slide_data.get('title', 'Thank You')}")
            if slide_data.get("subtitle"):
                lines.append(f"## {slide_data['subtitle']}")
        
        lines.append("")
    
    with open(output_path, "w") as f:
        f.write("\n".join(lines))
    
    return len(slides_data)


def main():
    parser = argparse.ArgumentParser(
        description="Create presentation slides from structured content",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  python slides.py --content slides.json --output presentation.pptx
  python slides.py --content slides.json --theme startup --output pitch.pptx
  python slides.py --content slides.json --format markdown --output slides.md

Themes: modern-dark, modern-light, corporate, startup, minimal
        """
    )
    
    parser.add_argument("--content", "-c", required=True,
                        help="JSON file with slide content")
    parser.add_argument("--output", "-o", required=True,
                        help="Output file path")
    parser.add_argument("--theme", "-t", default="modern-dark",
                        choices=list(THEMES.keys()),
                        help="Presentation theme (default: modern-dark)")
    parser.add_argument("--format", "-f", default="pptx",
                        choices=["pptx", "markdown", "md"],
                        help="Output format (default: pptx)")
    
    args = parser.parse_args()
    
    # Load content
    content_path = Path(args.content)
    if not content_path.exists():
        print(f"Error: Content file not found: {args.content}", file=sys.stderr)
        sys.exit(1)
    
    with open(content_path) as f:
        content = json.load(f)
    
    print("=" * 60)
    print("Slide Generation")
    print("=" * 60)
    print()
    
    # Generate output
    output_format = args.format.lower()
    if output_format in ["markdown", "md"]:
        num_slides = create_markdown(content, args.output)
    else:
        num_slides = create_pptx(content, args.theme, args.output)
    
    print(f"✅ Presentation created: {args.output}")
    print(f"   Slides: {num_slides}")
    print(f"   Theme: {args.theme}")
    print(f"   Format: {output_format}")


if __name__ == "__main__":
    main()
